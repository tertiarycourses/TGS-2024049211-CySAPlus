"""Branded Tertiary Infotech PPT design system (python-pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

SP = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(SP, "assets")

BLUE   = RGBColor(0x1F, 0x6F, 0xEB)
GREEN  = RGBColor(0x0F, 0xB9, 0x80)
INK    = RGBColor(0x16, 0x1B, 0x26)
GREY   = RGBColor(0x5B, 0x63, 0x72)
CARD   = RGBColor(0xF3, 0xF6, 0xFB)
CARDLN = RGBColor(0xE2, 0xE8, 0xF0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NAVY   = RGBColor(0x14, 0x2A, 0x54)
ACCENTS = [BLUE, GREEN, RGBColor(0x7C,0x3A,0xED), RGBColor(0xF5,0x9E,0x0B), RGBColor(0xE1,0x1D,0x48), RGBColor(0x0E,0xA5,0xE9)]

FONT = "Arial"
SW = Inches(13.333)
SH = Inches(7.5)
ML = Inches(0.62)
MR = Inches(0.62)
CW = SW - ML - MR

class Deck:
    def __init__(self, course_title, course_code, footer_copy="© 2026 Tertiary Infotech Academy Pte Ltd"):
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self.blank = self.prs.slide_layouts[6]
        self.course_title = course_title
        self.course_code = course_code
        self.footer_copy = footer_copy
        self.page = 0

    # ---------- primitives ----------
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _rect(self, s, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, shadow=False):
        sp = s.shapes.add_shape(shape, l, t, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line; sp.line.width = Pt(0.75)
        sp.shadow.inherit = False
        return sp

    def _text(self, s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              space_after=6, line_spacing=1.0, wrap=True):
        """runs: list of paragraphs; each paragraph = list of (text, size, bold, color, font)"""
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        for i, para in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(space_after); p.space_before = Pt(0)
            p.line_spacing = line_spacing
            for (txt, size, bold, color, *rest) in para:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(size); r.font.bold = bold
                r.font.name = rest[0] if rest else FONT
                r.font.color.rgb = color
        return tb

    def _footer(self, s):
        self._text(s, ML, SH - Inches(0.42), Inches(6), Inches(0.3),
                   [[(f"{self.course_title}  ·  {self.course_code}", 8.5, False, GREY)]])
        self._text(s, SW/2 - Inches(2.5), SH - Inches(0.42), Inches(5), Inches(0.3),
                   [[(self.footer_copy, 8.5, False, GREY)]], align=PP_ALIGN.CENTER)
        self._text(s, SW - MR - Inches(0.8), SH - Inches(0.42), Inches(0.8), Inches(0.3),
                   [[(str(self.page), 9, True, GREY)]], align=PP_ALIGN.RIGHT)

    def _header(self, s, eyebrow, title):
        # small blue vertical bar left of title
        self._rect(s, ML, Inches(0.45), Inches(0.09), Inches(0.72), BLUE)
        x = ML + Inches(0.26)
        self._text(s, x, Inches(0.42), CW - Inches(0.26), Inches(0.3),
                   [[(eyebrow.upper(), 13, True, BLUE)]])
        self._text(s, x, Inches(0.66), CW - Inches(0.26), Inches(0.7),
                   [[(title, 27, True, INK)]])
        # divider line
        ln = self._rect(s, ML, Inches(1.5), CW, Pt(1.4), CARDLN)

    # ---------- slide types ----------
    def title_slide(self, eyebrow, title, meta_lines, version_line):
        self.page += 1
        s = self._slide()
        self._rect(s, 0, 0, SW, Inches(0.12), BLUE)              # top bar
        self._rect(s, 0, SH - Inches(0.12), SW, Inches(0.12), GREEN)  # bottom bar
        # logos
        self._pic_fit(s, os.path.join(ASSETS, "tertiary_logo.png"), ML, Inches(0.55), Inches(1.15), Inches(1.15))
        self._pic_fit(s, os.path.join(ASSETS, "cysa_logo.png"), SW - MR - Inches(1.35), Inches(0.5), Inches(1.35), Inches(1.35))
        self._text(s, ML, Inches(2.55), CW, Inches(0.4), [[(eyebrow.upper(), 15, True, BLUE)]])
        self._text(s, ML, Inches(2.95), CW, Inches(1.6), [[(title, 46, True, INK)]], line_spacing=1.0)
        self._rect(s, ML, Inches(4.55), Inches(2.6), Inches(0.08), GREEN)
        yy = Inches(4.95)
        for ln in meta_lines:
            self._text(s, ML, yy, CW, Inches(0.35), [[(ln, 14.5, False, GREY)]])
            yy += Inches(0.36)
        self._text(s, ML, yy + Inches(0.05), CW, Inches(0.35), [[(version_line, 14.5, True, BLUE)]])
        self._text(s, ML, SH - Inches(0.55), CW, Inches(0.3),
                   [[(f"{self.footer_copy}. All rights reserved.  ·  www.tertiarycourses.com.sg", 8.5, False, GREY)]])
        return self.page

    def section_slide(self, eyebrow, title):
        self.page += 1
        s = self._slide()
        self._rect(s, 0, 0, Inches(0.22), SH, BLUE)  # left full bar
        self._rect(s, ML + Inches(0.1), Inches(3.1), Inches(0.12), Inches(1.35), GREEN)
        x = ML + Inches(0.42)
        self._text(s, x, Inches(3.15), CW, Inches(0.4), [[(eyebrow.upper(), 15, True, BLUE)]])
        self._text(s, x, Inches(3.55), CW, Inches(1.1), [[(title, 40, True, INK)]])
        self._footer(s)
        return self.page

    def content_slide(self, eyebrow, title, bullets):
        """bullets: list of (text, level) or (text, level, bold)"""
        self.page += 1
        s = self._slide()
        self._header(s, eyebrow, title)
        y = Inches(1.75)
        tb = s.shapes.add_textbox(ML, y, CW, SH - y - Inches(0.55))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for item in bullets:
            txt = item[0]; lvl = item[1] if len(item) > 1 else 0
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else INK
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.level = lvl
            p.space_after = Pt(7 if lvl == 0 else 3); p.line_spacing = 1.05
            if txt == "":
                continue
            # bullet marker
            rb = p.add_run(); rb.text = ("•  " if lvl == 0 else "–  ")
            rb.font.size = Pt(19 if lvl == 0 else 16); rb.font.bold = True
            rb.font.color.rgb = BLUE if lvl == 0 else GREEN; rb.font.name = FONT
            r = p.add_run(); r.text = txt
            r.font.size = Pt(19 if lvl == 0 else 16); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT
        self._footer(s)
        return self.page

    def cards_slide(self, eyebrow, title, items, cols=2):
        """items: list of (number_or_none, text)"""
        self.page += 1
        s = self._slide()
        self._header(s, eyebrow, title)
        n = len(items); rows = (n + cols - 1) // cols
        gx, gy = Inches(0.3), Inches(0.28)
        top = Inches(1.8)
        area_w = CW; area_h = SH - top - Inches(0.6)
        cw = int((area_w - gx * (cols - 1)) / cols)
        ch = int((area_h - gy * (rows - 1)) / rows)
        for i, (num, text) in enumerate(items):
            r = i // cols; c = i % cols
            l = ML + c * (cw + gx); t = top + r * (ch + gy)
            self._rect(s, l, t, cw, ch, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=CARDLN)
            acc = ACCENTS[i % len(ACCENTS)]
            self._rect(s, l, t, Inches(0.09), ch, acc, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tx = l + Inches(0.3)
            if num is not None:
                bd = Inches(0.62)
                self._rect(s, l + Inches(0.32), t + int(ch/2) - int(bd/2), bd, bd, acc, shape=MSO_SHAPE.OVAL)
                self._text(s, l + Inches(0.32), t + int(ch/2) - int(bd/2) - Inches(0.02), bd, bd,
                           [[(str(num), 18, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                tx = l + Inches(1.15)
            self._text(s, tx, t, l + cw - tx - Inches(0.2), ch,
                       [[(text, 16, False, INK)]], anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s)
        return self.page

    def steps_slide(self, eyebrow, title, steps):
        """steps: list of dict(label, desc, code)"""
        self.page += 1
        s = self._slide()
        self._header(s, eyebrow, title)
        y = Inches(1.72)
        tb = s.shapes.add_textbox(ML, y, CW, SH - y - Inches(0.55))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for st in steps:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.space_before = Pt(4); p.space_after = Pt(1); p.line_spacing = 1.03
            rb = p.add_run(); rb.text = "▍ "; rb.font.size = Pt(15); rb.font.bold = True; rb.font.color.rgb = BLUE; rb.font.name = FONT
            r = p.add_run(); r.text = st['label']; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = INK; r.font.name = FONT
            if st.get('desc'):
                pd = tf.add_paragraph(); pd.space_after = Pt(1); pd.line_spacing = 1.0
                pd.level = 1
                rd = pd.add_run(); rd.text = st['desc']; rd.font.size = Pt(12.5); rd.font.color.rgb = GREY; rd.font.name = FONT
            if st.get('code'):
                pc = tf.add_paragraph(); pc.space_after = Pt(5); pc.level = 1
                rc = pc.add_run(); rc.text = st['code']; rc.font.size = Pt(12); rc.font.name = "Consolas"; rc.font.color.rgb = RGBColor(0x0B,0x51,0x94)
        self._footer(s)
        return self.page

    def diagram_slide(self, eyebrow, title, img_files, caption=None):
        self.page += 1
        s = self._slide()
        self._header(s, eyebrow, title)
        files = [f for f in img_files if os.path.exists(f)]
        top = Inches(1.75); avail_w = CW; avail_h = SH - top - Inches(0.7)
        n = len(files)
        if n:
            gap = Inches(0.3)
            cellw = int((avail_w - gap * (n - 1)) / n)
            xs = ML
            total_w = 0; placed = []
            for f in files:
                iw, ih = Image.open(f).size; ar = iw / ih
                w = cellw; h = int(w / ar)
                if h > avail_h: h = int(avail_h); w = int(h * ar)
                placed.append((f, w, h))
                total_w += w
            total_w += gap * (n - 1)
            x = ML + int((CW - total_w) / 2)
            for f, w, h in placed:
                t = top + int((avail_h - h) / 2)
                s.shapes.add_picture(f, x, t, width=w, height=h)
                x += w + gap
        if caption:
            self._text(s, ML, SH - Inches(0.72), CW, Inches(0.3), [[(caption, 11, False, GREY)]], align=PP_ALIGN.CENTER)
        self._footer(s)
        return self.page

    def closing_slide(self, title, subtitle_lines):
        self.page += 1
        s = self._slide()
        self._rect(s, 0, 0, SW, SH, NAVY)
        self._rect(s, 0, SH - Inches(0.12), SW, Inches(0.12), GREEN)
        self._pic_fit(s, os.path.join(ASSETS, "tertiary_logo.png"), SW/2 - Inches(0.7), Inches(1.7), Inches(1.4), Inches(1.4))
        self._text(s, ML, Inches(3.4), CW, Inches(1.0), [[(title, 44, True, WHITE)]], align=PP_ALIGN.CENTER)
        yy = Inches(4.6)
        for ln in subtitle_lines:
            self._text(s, ML, yy, CW, Inches(0.4), [[(ln, 15, False, RGBColor(0xC7,0xD3,0xE8))]], align=PP_ALIGN.CENTER)
            yy += Inches(0.4)
        return self.page

    def _pic_fit(self, s, path, l, t, maxw, maxh):
        iw, ih = Image.open(path).size; ar = iw / ih
        w = maxw; h = int(w / ar)
        if h > maxh: h = maxh; w = int(h * ar)
        s.shapes.add_picture(path, l + int((maxw - w)/2), t + int((maxh - h)/2), width=w, height=h)

    def save(self, path):
        self.prs.save(path)
